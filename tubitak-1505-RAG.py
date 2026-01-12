import os
import time
import json
import hashlib
import tempfile
import requests
import torch
import re
from uuid import uuid4

import streamlit as st
from streamlit_lottie import st_lottie

from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_core.tools import tool
from hallucination_validator import HallucinationValidator

# --- SAYFA AYARLARI ---
st.set_page_config(page_title="Bilimp AI Asistan", layout="wide", page_icon="🤖")


# --- ANİMASYON YÜKLEME ---
def load_lottieurl(url: str):
    try:
        r = requests.get(url)
        if r.status_code != 200: return None
        return r.json()
    except:
        return None


# --- YARDIMCI FONKSİYON: METİN AKIŞI SİMÜLASYONU ---
# Tool kullanılmadığında hazır olan metni akışkan göstermek için
def stream_text_generator(text):
    for word in text.split(" "):
        yield word + " "
        time.sleep(0.05)


# --- GÖRSEL YÜKLEME EKRANI ---
if "app_loaded" not in st.session_state:
    loader_placeholder = st.empty()
    with loader_placeholder.container():
        st.markdown(
            """<style>.stApp {background-color: #0e1117;} .glowing-text {font-family: 'Source Code Pro', monospace; color: #00fbff; text-align: center; font-size: 2em; font-weight: bold; text-shadow: 0 0 10px #00fbff; animation: pulse 1.5s infinite;} @keyframes pulse { from {opacity: 0.8;} to {opacity: 1;} }</style>""",
            unsafe_allow_html=True)
        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            lottie_json = load_lottieurl("https://lottie.host/file/9860f43c-6232-4665-ba4f-557c669299b6.json")
            if lottie_json: st_lottie(lottie_json, height=250, key="loader", speed=1.5)

        status_text_placeholder = st.empty()
        loading_steps = ["🧠 Nöral Ağlar Yükleniyor...", "⚡ GPU Hızlandırma Aktif...",
                         "🛠️ Streaming (Akış) Modülü Başlatılıyor...", "🚀 Lütfen Bekleyiniz Sistem Hazırlanıyor..."]
        for step in loading_steps:
            status_text_placeholder.markdown(f'<p class="glowing-text">{step}</p>', unsafe_allow_html=True)
            time.sleep(0.5)

        # --- IMPORTLAR ---
        from qdrant_client import QdrantClient
        from qdrant_client.http import models as rest_models
        from qdrant_client.http.models import Distance, VectorParams, SparseVectorParams, Filter, FieldCondition, \
            MatchValue, MatchAny
        from langchain_qdrant import QdrantVectorStore, FastEmbedSparse, RetrievalMode
        from langchain_huggingface import HuggingFaceEmbeddings
        from langchain_google_genai import ChatGoogleGenerativeAI
        from langchain_ollama import ChatOllama
        from langchain_core.output_parsers import StrOutputParser
        import pymupdf4llm
        from markitdown import MarkItDown
        from pptx import Presentation
        from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter
        from langchain_core.documents import Document

    loader_placeholder.empty()
    st.session_state["app_loaded"] = True
else:
    from qdrant_client import QdrantClient
    from qdrant_client.http import models as rest_models
    from qdrant_client.http.models import Distance, VectorParams, SparseVectorParams, Filter, FieldCondition, \
        MatchValue, MatchAny
    from langchain_qdrant import QdrantVectorStore, FastEmbedSparse, RetrievalMode
    from langchain_huggingface import HuggingFaceEmbeddings
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_ollama import ChatOllama
    from langchain_core.output_parsers import StrOutputParser
    import pymupdf4llm
    from markitdown import MarkItDown
    from pptx import Presentation
    from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter
    from langchain_core.documents import Document

# ==============================================================================
# AYARLAR
# ==============================================================================
QDRANT_URL = "http://localhost:6333"
COLLECTION_NAME = "Tubitak_Dokumanlar_Hybrid"
EMBEDDING_MODEL_NAME = "ytu-ce-cosmos/turkish-e5-large"
REGISTRY_FILE = "belge_kayitlari.json"


# ==============================================================================
# YARDIMCI FONKSİYONLAR
# ==============================================================================
def calculate_md5(file_bytes):
    hash_md5 = hashlib.md5()
    hash_md5.update(file_bytes)
    return hash_md5.hexdigest()


def load_registry():
    if os.path.exists(REGISTRY_FILE):
        with open(REGISTRY_FILE, "r", encoding="utf-8") as f: return json.load(f)
    return {}


def save_registry(registry):
    with open(REGISTRY_FILE, "w", encoding="utf-8") as f: json.dump(registry, f, ensure_ascii=False, indent=4)


def delete_document_globally(filename):
    delete_by_source(filename)
    reg = load_registry()
    if filename in reg:
        del reg[filename]
        save_registry(reg)


def get_allowed_permissions(role):
    hierarchy = {"public": ["public"], "user": ["public", "user"], "management": ["public", "user", "management"],
                 "admin": ["public", "user", "management", "admin", "private"], "private": ["private"]}
    return hierarchy.get(role, ["public"])


def get_local_ollama_models():
    try:
        response = requests.get("http://localhost:11434/api/tags", timeout=1)
        if response.status_code == 200: return [m["name"] for m in response.json().get("models", [])]
    except:
        return []
    return []


# ==============================================================================
# CHUNKLAMA VE PARSE İŞLEMLERİ
# ==============================================================================
def etiketleri_generic_duzelt(text):
    lines = text.split('\n')
    new_lines = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") or (stripped.startswith("**") and stripped.endswith("**")):
            new_lines.append(line)
        else:
            new_lines.append(line)
    return '\n'.join(new_lines)


def process_pptx_native(file_path, source_name, permission):
    prs = Presentation(file_path)
    slides_chunks = []
    for i, slide in enumerate(prs.slides):
        content = []
        if slide.shapes.title and slide.shapes.title.text: content.append(f"# {slide.shapes.title.text.strip()}")
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame: content.append(shape.text.strip())
        full = "\n\n".join(content)
        if full.strip():
            doc = Document(page_content=full, metadata={"source": source_name, "chunk_no": i + 1, "file_type": "pptx",
                                                        "permission": permission})
            slides_chunks.append(doc)
    return slides_chunks


def process_text_based(file_path, source_name, chunk_size, chunk_overlap, permission):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        # 1. Markdown Dönüşümü
        if ext == ".pdf":
            text = pymupdf4llm.to_markdown(file_path, write_images=False)
        else:
            md = MarkItDown()
            result = md.convert(file_path)
            text = result.text_content

        # Temizlik
        clean = etiketleri_generic_duzelt(text)

        # 2. Başlıklara Göre Bölme
        headers_to_split_on = [
            ("#", "Main"),
            ("##", "Sub"),
            ("###", "Sub2"),
            ("####", "Sub3")
        ]

        splitter = MarkdownHeaderTextSplitter(
            headers_to_split_on=headers_to_split_on,
            strip_headers=True
        )
        md_docs = splitter.split_text(clean)

        # --- AKILLI BİRLEŞTİRME 2.0 (SAFE MERGE) ---
        merged_docs = []
        temp_doc = None

        for doc in md_docs:
            if not doc.page_content.strip():
                continue

            # Context (Bağlam) bilgisini hazırla
            header_path = " > ".join([doc.metadata.get(h[1]) for h in headers_to_split_on if doc.metadata.get(h[1])])
            if header_path:
                doc.page_content = f"**BAĞLAM:** {header_path}\n\n{doc.page_content}"

            # --- MANTIK BAŞLIYOR ---

            # Eğer elimizde bekleyen "yetim" bir parça varsa:
            if temp_doc:
                # GÜVENLİK KONTROLÜ: Şimdiki parça da çok kısaysa (başka bir başlık olabilir), birleştirme yapma!
                # Çünkü "Yemek Listesi" başlığı ile "Servis Saatleri" başlığını birleştirmek istemeyiz.
                if len(doc.page_content) < 100 and "|" not in doc.page_content:
                    # Bekleyeni olduğu gibi kaydet, çünkü arkasından gelen de içerik değilmiş.
                    merged_docs.append(temp_doc)
                    temp_doc = doc  # Şimdikini yeni bekleyen yap
                else:
                    # Şimdiki parça dolu bir içerik (Tablo veya Uzun Metin). Birleştir!
                    # Önceki kısa başlık + Yeni Satır + Şimdiki İçerik
                    new_content = f"{temp_doc.page_content}\n\n{doc.page_content}"
                    doc.page_content = new_content
                    # Metadata'yı koru (genelde aynı başlık altındadırlar)
                    merged_docs.append(doc)
                    temp_doc = None  # Bekleyen kutusunu boşalt

            else:
                # Elimizde bekleyen yok. Peki bu parça beklemeye alınmalı mı?
                # Kural: 250 karakterden kısaysa VE içinde Tablo yoksa -> Potansiyel Yetim Başlık
                if len(doc.page_content) < 250 and "|" not in doc.page_content:
                    temp_doc = doc
                else:
                    # Parça zaten büyük veya tablo, direkt ekle.
                    merged_docs.append(doc)

        # Döngü bittiğinde elde kalan son parça varsa onu da ekle (Unutma!)
        if temp_doc:
            merged_docs.append(temp_doc)
        # -------------------------------------------------------------

        # 3. Recursive Splitter (Çok büyükleri bölmek için)
        rec_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", " ", ""]
        )

        final_docs = []
        for doc in merged_docs:
            doc.metadata.update({
                "source": source_name,
                "file_type": ext.replace(".", ""),
                "permission": permission
            })
            chunks = rec_splitter.split_documents([doc])
            final_docs.extend(chunks)

        return final_docs

    except Exception as e:
        st.error(f"Hata: {e}")
        return []


# ==============================================================================
# QDRANT VE EMBEDDING MODELLERİ
# ==============================================================================
@st.cache_resource
def get_dense_embeddings():
    device = "cuda" if torch.cuda.is_available() else "cpu"
    return HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL_NAME, model_kwargs={"device": device},
                                 encode_kwargs={"normalize_embeddings": True})


@st.cache_resource
def get_sparse_embeddings():
    return FastEmbedSparse(model_name="Qdrant/bm25")


@st.cache_resource
def get_qdrant_client():
    return QdrantClient(url=QDRANT_URL, check_compatibility=False)


def init_collection():
    client = get_qdrant_client()
    if not client.collection_exists(COLLECTION_NAME):
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config={"content": VectorParams(size=1024, distance=Distance.COSINE)},
            sparse_vectors_config={"sparse": SparseVectorParams()}
        )


def add_documents_to_qdrant(documents):
    client = get_qdrant_client()
    dense_emb = get_dense_embeddings()
    sparse_emb = get_sparse_embeddings()
    vector_store = QdrantVectorStore(client=client, collection_name=COLLECTION_NAME, embedding=dense_emb,
                                     vector_name="content", sparse_embedding=sparse_emb, sparse_vector_name="sparse",
                                     retrieval_mode=RetrievalMode.HYBRID)
    ids = [str(uuid4()) for _ in documents]
    vector_store.add_documents(documents=documents, ids=ids)


def delete_by_source(source_name):
    client = get_qdrant_client()
    if client.collection_exists(COLLECTION_NAME):
        client.delete(collection_name=COLLECTION_NAME, points_selector=Filter(
            must=[FieldCondition(key="metadata.source", match=MatchValue(value=source_name))]))


# ==============================================================================
# ARAYÜZ
# ==============================================================================
with st.sidebar:
    try:
        st.image("bilimp_logo.png", width="stretch")
    except:
        st.warning("Logo Yok")

    st.markdown("### 🛠️ Sistem Ayarları")
    if "last_role" not in st.session_state: st.session_state.last_role = "admin"
    current_user_role = st.selectbox("👤 Kullanıcı Rolü", ["public", "user", "management", "admin", "private"], index=3)
    if current_user_role != st.session_state.last_role:
        st.session_state.messages = []
        st.session_state.last_role = current_user_role
        st.rerun()
    with st.expander("ℹ️ Yetki Detayı"):
        st.code(get_allowed_permissions(current_user_role))

    st.divider()
    st.markdown("### 🧠 Yapay Zeka Motoru")
    gemini_models_map = {"Gemini 2.5 Flash (Hızlı)": "gemini-2.5-flash",
                         "Gemini 3.0 Flash (Akıllı + Hızlı)": "gemini-3-flash-preview"}
    ollama_list = get_local_ollama_models()
    model_options = list(gemini_models_map.keys())
    if ollama_list:
        model_options.extend([f"Ollama: {m}" for m in ollama_list])
    else:
        model_options.append("Ollama (Model Yok)")

    selected_option = st.selectbox("Model Seçimi", model_options)
    llm_model_id, llm_type = None, "ollama"
    if "Gemini" in selected_option:
        llm_type = "gemini";
        llm_model_id = gemini_models_map[selected_option]
    elif "Ollama" in selected_option:
        llm_type = "ollama";
        llm_model_id = selected_option.split(": ")[1]

    api_key = ""
    if llm_type == "gemini": api_key = st.text_input("🔑 Google API Key", type="password")

    st.divider()
    st.markdown("### 🎛️ İnce Ayarlar")
    temperature = st.slider("Yaratıcılık", 0.0, 1.0, 0.1, step=0.1)
    top_k = st.number_input("Bağlam (Chunk)", 1, 20, 5)
    score_threshold = st.slider("Benzerlik Eşiği", 0.0, 0.9, 0.70, step=0.05)
    with st.expander("📄 Chunk Parametreleri"):
        c_size = st.number_input("Boyut", 500, 5000, 2500)
        c_over = st.number_input("Örtüşme", 0, 1000, 200)

st.header("📄 Bilimp Doküman Asistanı (Streaming Agent)")
t1, t2 = st.tabs(["📂 **Belge Yönetimi**", "💬 **Akıllı Sohbet**"])

# --- TAB 1: BELGE YÖNETİMİ ---
with t1:
    col_upload, col_list = st.columns([1, 1], gap="large")
    with col_upload:
        st.markdown("#### ⬆️ Belge Yükle")
        up_file = st.file_uploader("Dosyayı buraya sürükleyin", type=["pdf", "docx", "xlsx", "pptx"],
                                   label_visibility="collapsed")
        if up_file:
            bytes_data = up_file.getvalue();
            f_name = up_file.name;
            curr_md5 = calculate_md5(bytes_data);
            reg = load_registry()
            file_exists = False
            if f_name in reg:
                stored = reg[f_name]
                if isinstance(stored, dict) and stored["hash"] == curr_md5:
                    file_exists = True
                elif stored == curr_md5:
                    file_exists = True

            if file_exists:
                st.warning(f"⚠️ **{f_name}** zaten mevcut.")
            else:
                st.success(f"✅ **{f_name}** analize hazır.")

            if st.button("🚀 Sisteme Entegre Et", type="primary"):
                with st.status("İşleniyor...", expanded=True) as s:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(f_name)[1]) as tmp:
                        tmp.write(bytes_data);
                        tmp_path = tmp.name
                    init_collection();
                    delete_by_source(f_name)
                    chunks = []
                    if f_name.endswith(".pptx"):
                        chunks = process_pptx_native(tmp_path, f_name, current_user_role)
                    else:
                        chunks = process_text_based(tmp_path, f_name, c_size, c_over, current_user_role)
                    if chunks:
                        add_documents_to_qdrant(chunks)
                        reg[f_name] = {"hash": curr_md5, "permission": current_user_role}
                        save_registry(reg)
                        s.update(label="Tamamlandı!", state="complete", expanded=False)
                        st.toast("Başarılı!", icon="🎉");
                        time.sleep(1);
                        st.rerun()
                    else:
                        s.update(label="Hata", state="error");
                        st.error("Ayrıştırılamadı.")
                    os.unlink(tmp_path)

    with col_list:
        st.markdown("#### 🗂️ Sistemdeki Belgeler")
        current_reg = load_registry();
        allowed_view_perms = get_allowed_permissions(current_user_role);
        visible_files = []
        for fname, fdata in current_reg.items():
            perm = fdata.get("permission", "public") if isinstance(fdata, dict) else "public"
            if perm in allowed_view_perms: visible_files.append((fname, perm))
        if not visible_files:
            st.info("Görüntülenecek belge yok.")
        else:
            for fname, perm in visible_files:
                c1, c2 = st.columns([0.8, 0.2])
                with c1:
                    st.markdown(
                        f"""<div style="padding:10px; background:#161b22; border-radius:8px; margin-bottom:5px; border:1px solid #30363d;"><span style="color:white; font-weight:600;">📄 {fname}</span><span style="background:#238636; color:white; padding:2px 8px; border-radius:4px; font-size:0.8em; margin-left:10px;">{perm}</span></div>""",
                        unsafe_allow_html=True)
                with c2:
                    if st.button("🗑️", key=f"del_{fname}"): delete_document_globally(fname); st.rerun()

# --- TAB 2: SOHBET (STREAMING) ---
# --- TAB 2: SOHBET (STREAMING) ---
with t2:
    # Memory için yardımcı fonksiyon
    def get_formatted_history(messages, max_pairs=5):
        """
        Mesaj geçmişini LangChain formatına çevirir.
        max_pairs: Maksimum user-assistant çifti sayısı
        """
        history = []
        # Tüm mesajları al
        all_msgs = messages.copy()

        # Son N çifti almak için (her çift = 1 user + 1 assistant)
        # En fazla max_pairs * 2 mesaj al
        recent = all_msgs[-(max_pairs * 2):]

        for msg in recent:
            content = msg.get("content", "")
            if not content or content.strip() == "":
                continue

            if msg["role"] == "user":
                history.append(HumanMessage(content=content))
            elif msg["role"] == "assistant":
                history.append(AIMessage(content=content))

        return history


    # Mesaj geçmişini başlat
    if "messages" not in st.session_state:
        st.session_state.messages = []

    # Geçmiş mesajları render et
    for m in st.session_state.messages:
        with st.chat_message(m["role"]):
            st.markdown(m["content"])
            if m["role"] == "assistant" and "sources" in m and m["sources"]:
                with st.expander(f"🔍 Referans Kaynaklar ({len(m['sources'])})"):
                    for i, doc in enumerate(m['sources']):
                        score_val = doc.metadata.get("score", 0.0)
                        st.markdown(f"**#{i + 1}** | 📂 `{doc.metadata.get('source')}` | 📊 Skor: `{score_val:.4f}`")
                        st.caption(doc.page_content)
                        st.divider()

    if prompt := st.chat_input("Sorunuzu buraya yazın..."):
        # Kullanıcı mesajını ÖNCE ekle
        st.session_state.messages.append({"role": "user", "content": prompt})

        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            client = get_qdrant_client()
            if not client.collection_exists(COLLECTION_NAME):
                st.error("Veritabanı boş.")
            else:
                ready = True
                llm = None
                if llm_type == "gemini":
                    if not api_key:
                        st.error("API Key Eksik!")
                        ready = False
                    else:
                        llm = ChatGoogleGenerativeAI(
                            model=llm_model_id,
                            google_api_key=api_key,
                            temperature=temperature
                        )
                elif llm_type == "ollama":
                    if "Yok" in selected_option:
                        st.error("Model Yok!")
                        ready = False
                    else:
                        llm = ChatOllama(model=llm_model_id, temperature=temperature)

                if ready and llm:
                    try:
                        @tool
                        def bilimp_knowledge_base(query: str):
                            """
                            Bilimp AI Asistanı'nın şirket içi bilgi bankasında arama yapmasını sağlar.
                            """
                            try:
                                dense_emb = get_dense_embeddings()
                                sparse_emb = get_sparse_embeddings()
                                vector_store = QdrantVectorStore(
                                    client=client,
                                    collection_name=COLLECTION_NAME,
                                    embedding=dense_emb,
                                    vector_name="content",
                                    sparse_embedding=sparse_emb,
                                    sparse_vector_name="sparse",
                                    retrieval_mode=RetrievalMode.HYBRID
                                )
                                
                                allowed_perms = get_allowed_permissions(current_user_role)
                                perm_filter = rest_models.Filter(must=[
                                    rest_models.FieldCondition(
                                        key="metadata.permission",
                                        match=rest_models.MatchAny(any=allowed_perms)
                                    )
                                ])
                                
                                results = vector_store.similarity_search_with_score(
                                    query, k=5, filter=perm_filter
                                )
                                
                                high_quality_docs = [
                                    doc for doc, score in results 
                                    if score >= 0.70
                                ]
                                
                                if not high_quality_docs:
                                    return "ARAŞTIRMA_SONUCU: Bu konuda belgelerimde yeterli kalitede bilgi bulunamadı."
                                
                                if len(high_quality_docs) < 2:
                                    return "ARAŞTIRMA_SONUCU: Konu hakkında çok az bilgi var, lütfen daha spesifik soru sorun."
                                
                                context = "\n\n".join([doc.page_content for doc in high_quality_docs])
                                return f"ARAŞTIRMA_SONUCU: {context}"
                                
                            except Exception as e:
                                return f"ARAŞTIRMA_SONUCU: Teknik hata - {str(e)}"


                        llm_with_tools = llm.bind_tools([bilimp_knowledge_base])

                        # ---------------------------------------------------------
                        # 2. GEÇMİŞİ DÜZGÜN FORMATTA HAZIRLA (KRİTİK DEĞİŞİKLİK)
                        # ---------------------------------------------------------
                        # SON mesajı (şu anki prompt) HARİÇ tutarak geçmişi al
                        # Çünkü prompt zaten ayrıca ekleniyor
                        history_messages = st.session_state.messages[:-1]  # Son mesaj hariç
                        history_langchain_format = get_formatted_history(history_messages, max_pairs=5)

                        # ---------------------------------------------------------
                        # 3. SİSTEM PROMPTU
                        # ---------------------------------------------------------
                        identity_section = """
                                Sen profesyonel, yardımsever ve kurumsal bir asistansın.
                                KİMLİĞİN:
                                - Adın: **Bilimp AI Asistanı**.
                                - Görevin: Çalışanlara şirket içi dökümanlar, yönetmelikler ve prosedürler hakkında bilgi sağlamak.
                                
                                YETENEKLERİN VE HAFIZA:
                                - Güçlü bir hafızan var. Sohbet geçmişindeki TÜM mesajları (hem kullanıcının sorularını HEM DE kendi verdiğin cevapları) hatırlarsın.
                                - Kullanıcı "Önceki soruma ne cevap verdin?", "Az önce ne dedin?", "Bir önceki cevabın neydi?" gibi sorular sorarsa, sohbet geçmişine bakarak KENDİ VERDİĞİN CEVAPLARI söyle.
                                - Kullanıcı "Ne sormuştum?" derse, onun önceki sorularını hatırla.
                                
                                DAVRANIŞ KURALLARI:
                                1. Eğer kullanıcı "Kimsin?" derse kendini tanıt.
                                2. Başka bir model olduğunu ASLA SÖYLEME.
                                3. Kullanıcıya her zaman nazik ve "siz" diliyle hitap et.
                                4. Hafıza soruları için TOOL KULLANMA, direkt sohbet geçmişinden cevapla.
                                """

                        router_section = """
                                GÖREVİN:
                                Gelen soruyu ve sohbet geçmişini analiz edip 'bilimp_knowledge_base' aracını kullanıp kullanmayacağına karar ver.
                                
                                KARAR MANTIĞI:
                                1. **Veri İsteği:** Şirket verisi, sayı, kural soruluyorsa -> TOOL KULLAN.
                                2. **Takip Sorusu:** "Peki kaç tane?", "Bunun fiyatı ne?" gibi önceki konunun devamıysa -> TOOL KULLAN.
                                3. **HAFIZA SORULARI (KRİTİK):**
                                   - "Önceki cevabın neydi?", "Ne demiştin?", "Az önce ne söyledin?" -> TOOL KULLANMA, sohbet geçmişinden cevapla.
                                   - "Ne sormuştum?", "Önceki sorum neydi?" -> TOOL KULLANMA, sohbet geçmişinden cevapla.
                                4. **Sohbet:** "Merhaba", "Nasılsın" -> TOOL KULLANMA.
                                """

                        full_system_prompt = identity_section + "\n\n" + router_section

                        # ---------------------------------------------------------
                        # 4. MESAJLARI OLUŞTUR
                        # ---------------------------------------------------------
                        input_msgs = [
                                         SystemMessage(content=full_system_prompt)
                                     ] + history_langchain_format + [
                                         HumanMessage(content=prompt)
                                     ]

                        # DEBUG: Geçmişi kontrol et (geliştirme aşamasında kullan)
                        # st.write(f"📜 Geçmiş mesaj sayısı: {len(history_langchain_format)}")
                        # for i, msg in enumerate(history_langchain_format):
                        #     st.write(f"{i}: {type(msg).__name__} - {msg.content[:50]}...")

                        # Model Karar Veriyor
                        ai_msg = llm_with_tools.invoke(input_msgs)

                        # Değişkenleri sıfırla
                        final_response = ""
                        retrieved_docs = []

                        # ---------------------------------------------------------
                        # 5. DURUMA GÖRE CEVAPLAMA
                        # ---------------------------------------------------------
                        if ai_msg.tool_calls:
                            # DURUM A: RAG GEREKLİ
                            with st.status("📚 Bilgi Bankası Taranıyor...", expanded=True) as s:
                                dense_emb = get_dense_embeddings()
                                sparse_emb = get_sparse_embeddings()
                                vector_store = QdrantVectorStore(
                                    client=client,
                                    collection_name=COLLECTION_NAME,
                                    embedding=dense_emb,
                                    vector_name="content",
                                    sparse_embedding=sparse_emb,
                                    sparse_vector_name="sparse",
                                    retrieval_mode=RetrievalMode.HYBRID
                                )
                                allowed_perms = get_allowed_permissions(current_user_role)
                                perm_filter = rest_models.Filter(must=[
                                    rest_models.FieldCondition(
                                        key="metadata.permission",
                                        match=rest_models.MatchAny(any=allowed_perms)
                                    )
                                ])

                                results = vector_store.similarity_search_with_score(prompt, k=top_k, filter=perm_filter)
                                for doc, score in results:
                                    if score >= score_threshold:
                                        doc.metadata["score"] = score
                                        retrieved_docs.append(doc)

                                if not retrieved_docs:
                                    s.update(label="Bilgi Bulunamadı", state="error", expanded=False)
                                    final_response = "Bu konuda belgelerimde yeterli kalitede bilgi bulunamadı. Lütfen sorunuzu farklı şekilde ifade edin."
                                    st.error("❌ " + final_response)
                                    st.session_state.messages.append({
                                        "role": "assistant",
                                        "content": final_response,
                                        "sources": []
                                    })
                                    st.stop()
                                
                                if len(retrieved_docs) < 2:
                                    s.update(label="Yetersiz Bilgi", state="warning", expanded=False)
                                    final_response = "Bu konuda çok az bilgi var. Lütfen daha spesifik bir soru sorun."
                                    st.warning("⚠️ " + final_response)
                                    st.session_state.messages.append({
                                        "role": "assistant",
                                        "content": final_response,
                                        "sources": retrieved_docs
                                    })
                                    st.stop()

                                context_str = "\n\n".join([d.page_content for d in retrieved_docs])
                                s.update(label="Bilgiler Getirildi!", state="complete", expanded=False)

                            rag_system_prompt = f"""
SİSTEM TALİMATI: Sen TÜBİTAK 1505 doküman uzmanısın.

BULUNAN DÖKÜMANLAR:
{context_str}

KATÎ KURALLAR:
1. SADECE yukarıdaki belgelerden cevap ver.
2. Belgeler soruyu tam cevaplamıyorsa bunu açıkça belirt.
3. Belirsizlik varsa "Mevcut belgeler bu konuda net bilgi içermiyor" de.
4. ASLA tahmin yapma, spekülasyon etme veya kendi bilgini ekleme.
5. Her cevabın sonunda hangi belgeden aldığını belirt.
6. Cevabın TAMAMEN Türkçe olmalıdır.

UYARI: Yukarıdaki belgeler soruyu cevaplamak için yetersizse bunu kullanıcıya söyle.
"""
                            st.markdown("📚 **Dökümanlardan Yanıtlanıyor:**")

                            # RAG için de geçmişi ekle
                            rag_messages = [
                                               SystemMessage(content=rag_system_prompt)
                                           ] + history_langchain_format + [
                                               HumanMessage(content=prompt)
                                           ]

                            stream_generator = llm.stream(rag_messages)
                            final_response = st.write_stream(stream_generator)
                            
                            is_valid, validation_msg = HallucinationValidator.validate_response(
                                prompt, final_response, retrieved_docs
                            )
                            
                            if not is_valid:
                                st.warning(f"⚠️ Kalite Uyarısı: {validation_msg}")
                                final_response = "Bu konuda belgelerimde net bilgi bulamadım. Lütfen sorunuzu farklı şekilde ifade edin."
                                st.error(final_response)

                        else:
                            # DURUM B: SOHBET (Tool Yok)
                            raw_content = ai_msg.content
                            content_text = ""

                            if isinstance(raw_content, str):
                                content_text = raw_content
                            elif isinstance(raw_content, list):
                                for item in raw_content:
                                    if isinstance(item, list):
                                        for sub_item in item:
                                            if isinstance(sub_item, dict):
                                                content_text += sub_item.get("text", "")
                                    elif isinstance(item, dict):
                                        content_text += item.get("text", "")
                                    elif isinstance(item, str):
                                        content_text += item
                            else:
                                content_text = str(raw_content)

                            st.markdown("💬 **Sohbet Modu:**")
                            final_response = st.write_stream(stream_text_generator(content_text))

                        # ---------------------------------------------------------
                        # 6. GEÇMİŞE KAYDET (KRİTİK!)
                        # ---------------------------------------------------------
                        if retrieved_docs:
                            with st.expander(f"🔍 Referans Kaynaklar ({len(retrieved_docs)})"):
                                for i, doc in enumerate(retrieved_docs):
                                    score_val = doc.metadata.get("score", 0.0)
                                    st.markdown(
                                        f"**#{i + 1}** | 📂 `{doc.metadata.get('source')}` | 📊 Skor: `{score_val:.4f}`")
                                    st.caption(doc.page_content)

                        # Assistant cevabını kaydet
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": final_response,  # Tam cevap metni
                            "sources": retrieved_docs
                        })

                    except Exception as e:
                        error_msg = str(e)
                        if "429" in error_msg:
                            st.error("⚠️ API Kotası Doldu.")
                        else:
                            st.error(f"Hata: {e}")
                        # Hata durumunda da boş mesaj ekleme
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": f"Bir hata oluştu: {error_msg}",
                            "sources": []
                        })